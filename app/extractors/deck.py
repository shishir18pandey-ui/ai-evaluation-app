
import logging
from pathlib import Path

from pypdf import PdfReader
from pptx import Presentation

from app.core.models import Evidence
from app.core import config
from app.llm import llm_json

logger = logging.getLogger(__name__)


CLAIM_EXTRACTION_SYSTEM = """You extract factual and technical claims from a multi-slide deck.

Input is the full deck text, with each slide marked like "=== SLIDE N ===".

A CLAIM is any concrete assertion about:
- what the product does (features)
- how it works (architecture, tech stack)
- what it has achieved (metrics, users, results)
- what problem it solves
- what is planned for the future (mark as "future" category)

NOT a claim: marketing fluff, taglines, team bios, generic statements.

Return JSON with ONE key "claims" containing an array. Each item:
{
  "slide": N,
  "claim": "short factual statement",
  "category": "feature|architecture|metric|problem|future|other",
  "verbatim_quote": "exact text from slide",
  "confidence": 0.0-1.0
}

Return {"claims": []} if no real claims exist.
Process ALL slides — do not stop early.
"""


def _read_pdf(path: Path) -> list[tuple[int, str]]:
    reader = PdfReader(str(path))
    return [(i + 1, page.extract_text() or "") for i, page in enumerate(reader.pages)]


def _read_pptx(path: Path) -> list[tuple[int, str]]:
    prs = Presentation(str(path))
    slides = []
    for i, slide in enumerate(prs.slides):
        parts = [shape.text for shape in slide.shapes if hasattr(shape, "text") and shape.text]
        slides.append((i + 1, "\n".join(parts)))
    return slides


def extract_from_deck(file_path: str | Path) -> list[Evidence]:
    """Main entry — returns list of Evidence. ONE LLM call for the whole deck."""
    path = Path(file_path)
    if not path.exists():
        logger.warning("Deck not found: %s", path)
        return []

    ext = path.suffix.lower()
    if ext == ".pdf":
        slides = _read_pdf(path)
    elif ext in (".pptx", ".ppt"):
        slides = _read_pptx(path)
    else:
        logger.error("Unsupported deck format: %s", ext)
        return []

    slides = slides[: config.MAX_SLIDES_TO_PROCESS]
    if not slides:
        return []

    # Build single labelled blob for ONE LLM call
    parts = []
    for slide_num, text in slides:
        clean = (text or "").strip()
        if not clean:
            continue
        parts.append(f"=== SLIDE {slide_num} ===\n{clean}")
    if not parts:
        return []

    deck_blob = "\n\n".join(parts)
    # Truncate aggressively if huge — protect token limit
    if len(deck_blob) > 20000:
        deck_blob = deck_blob[:20000] + "\n\n[... deck truncated ...]"

    logger.info("Extracting claims from %d slides in ONE call", len(slides))

    try:
        result = llm_json(
            CLAIM_EXTRACTION_SYSTEM,
            f"DECK:\n\n{deck_blob}\n\nExtract all claims.",
        )
    except Exception as e:
        logger.error("Deck extraction failed: %s", e)
        return []

    claims = result.get("claims", []) if isinstance(result, dict) else []
    evidences = []
    for c in claims:
        if not isinstance(c, dict) or "claim" not in c:
            continue
        slide_num = c.get("slide", 0)
        try:
            slide_num = int(slide_num)
        except (ValueError, TypeError):
            slide_num = 0
        evidences.append(Evidence(
            source_type="deck",
            source_id=f"slide_{slide_num}" if slide_num else "slide_unknown",
            claim_or_fact=c.get("claim", ""),
            evidence_text=c.get("verbatim_quote", "") or c.get("claim", ""),
            confidence=float(c.get("confidence", 0.7)),
            metadata={
                "slide_num": slide_num,
                "category": c.get("category", "other"),
            },
        ))

    logger.info("Extracted %d claims from deck", len(evidences))
    return evidences